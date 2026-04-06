import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # 0 = Title Slide, 1 = Title and Content
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]
    
    # 1. Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "The Spiritual Journey of Holy Week"
    subtitle.text = "Through Hymns & Rites"

    # 2. Entering a Different Kind of Time
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "1) Entering a Different Kind of Time"
    tf = slide.placeholders[1].text_frame
    tf.text = "Holy Week is not a normal week"
    p = tf.add_paragraph()
    p.text = "The Church invites us to slow down, step out of ordinary time, and walk with Christ"
    p = tf.add_paragraph()
    p.text = "We are not remembering events — we are walking with Christ"
    p = tf.add_paragraph()
    p.text = "Key question: Why is everything so long, repetitive, and musical?"
    
    # 3. Why This Week Exists: The Meaning of the Cross
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "2) Why This Week Exists: The Meaning of the Cross"
    tf = slide.placeholders[1].text_frame
    tf.text = "Christ entered our pain, fear, and death"
    p = tf.add_paragraph()
    p.text = "The Passion is divine love, not tragedy"
    p = tf.add_paragraph()
    p.text = "The Cross is where justice meets mercy"

    # 4. The Structure of the Journey
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "3) The Structure of the Journey"
    tf = slide.placeholders[1].text_frame
    tf.text = "Palm Sunday → rejection"
    p = tf.add_paragraph()
    p.text = "Monday–Wednesday → confrontation"
    p = tf.add_paragraph()
    p.text = "Thursday → love and betrayal"
    p = tf.add_paragraph()
    p.text = "Friday → sacrifice"
    p = tf.add_paragraph()
    p.text = "Saturday → waiting"
    p = tf.add_paragraph()
    p.text = "Resurrection → victory"
    p = tf.add_paragraph()
    p.text = "Framework: Readings, Rites, Hymns"
    
    # 5. The key triad of the Holy Week experience
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "The Key Triad of the Holy Week Experience"
    tf = slide.placeholders[1].text_frame
    tf.text = "Readings, Rites, Hymns"
    p = tf.add_paragraph()
    p.text = "Hymns are theology, prayer, and experience"
    p = tf.add_paragraph()
    p.text = "They express what words cannot"
    p = tf.add_paragraph()
    p.text = "The Church teaches theology through melody"

    # 6. Entering One Hour of Holy Week & Prophecies
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "5) Entering One Hour of Holy Week"
    tf = slide.placeholders[1].text_frame
    tf.text = "Introduction to the prophecies is an opening hymn to set a slow, sorrowful tone"
    p = tf.add_paragraph()
    p.text = "The Prophecies (Old Testament Readings)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Drawn from: Isaiah, Jeremiah, Torah, Minor prophets"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Theological meaning: The Cross was not accidental. It was foretold, prepared, and awaited"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Spiritual purpose: The Church wants you to see Christ everywhere in Scripture and realize that salvation is one continuous story"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "💡 Key insight: The Cross is not a moment… it is the fulfillment of all history"
    p.level = 1

    # 7. Homily
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Homily"
    tf = slide.placeholders[1].text_frame
    tf.text = "The precious teachings of the church fathers that reflect on the events and purpose of this hour or day."
    p = tf.add_paragraph()
    p.text = "Liturgical Wisdom: Before the homily the church chants a calm and beautiful tune as an introduction"
    p = tf.add_paragraph()
    p.text = "Why is it placed here? Because the readings are long and dense. The mind becomes saturated, tired, full."
    p = tf.add_paragraph()
    p.text = "Purpose of the calm melody: Not filler — but spiritual recovery. It allows digestion of the Word, emotional settling, preparation for deeper meaning."
    p = tf.add_paragraph()
    p.text = "💡 Key insight: The Church understands the human soul and builds rest into worship"

    # 8. Paschal Praise (Thine is the power)
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Paschal Praise (Thine is the power)"
    tf = slide.placeholders[1].text_frame
    tf.text = "Theological meaning: Declares Christ’s kingship and victory"
    p = tf.add_paragraph()
    p.text = "Paradox: We are in the middle of suffering, betrayal, crucifixion... Yet we say: “Yours is the power and the glory”"
    p = tf.add_paragraph()
    p.text = "Spiritual meaning: The Church sees the Cross not as defeat but as enthronement. The Cross = love + victory, not loss"
    p = tf.add_paragraph()
    p.text = "💡 Key insight: We proclaim victory before we see it"
    
    # 9. Repetition (Spiritual Formation)
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Repetition (Spiritual Formation)"
    tf = slide.placeholders[1].text_frame
    tf.text = "The Paschal Praise is repeated many times"
    p = tf.add_paragraph()
    p.text = "Why repetition? To move truth from mind → heart"
    p = tf.add_paragraph()
    p.text = "Deeper spiritual role: Creates rhythm, stillness, internal participation. Hymns express what words cannot."
    p = tf.add_paragraph()
    p.text = "💡 Key insight: Repetition turns knowledge into experience. Repetition engraves truth in the heart."

    # 10. Psalm in the long tune
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Psalm in the Long Tune"
    tf = slide.placeholders[1].text_frame
    tf.text = "Spiritual meaning: The Psalm is deeply emotional and personal"
    p = tf.add_paragraph()
    p.text = "So the Church stretches it, gives space to feel it"
    p = tf.add_paragraph()
    p.text = "Effect: Moves from communal reading → personal prayer"
    p = tf.add_paragraph()
    p.text = "💡 Key insight: The Psalm is where you enter the story"

    # 11. The Gospel
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "The Gospel"
    tf = slide.placeholders[1].text_frame
    tf.text = "The climax of the hour"
    p = tf.add_paragraph()
    p.text = "Meaning: Now Christ speaks directly; the Passion unfolds"
    p = tf.add_paragraph()
    p.text = "Spiritual significance: Everything before was preparation / Now comes encounter"
    p = tf.add_paragraph()
    p.text = "Connection to structure: Readings → prepare / Hymns → soften / Gospel → reveal"
    p = tf.add_paragraph()
    p.text = "💡 Key insight: The Gospel is not information - It is presence"

    # 12. Exposition (Explanation of the Hour)
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Exposition (Explanation of the Hour)"
    tf = slide.placeholders[1].text_frame
    tf.text = "A poetic, often symbolic explanation"
    p = tf.add_paragraph()
    p.text = "Function: Interprets events, symbols and spiritual meaning"
    p = tf.add_paragraph()
    p.text = "Style: Not academic. Often poetic, symbolic, emotional."
    p = tf.add_paragraph()
    p.text = "Spiritual role: Helps the believer locate themselves in the story, understand the deeper meaning"
    p = tf.add_paragraph()
    p.text = "💡 Key insight: The exposition teaches us how to see spiritually"

    # 13. The Hour as a whole
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "The Hour as a Whole"
    tf = slide.placeholders[1].text_frame
    tf.text = "Why is it long? Because each element prepares, deepens, transforms"
    p = tf.add_paragraph()
    p.text = "Holy Week is meant for full immersion and devotion"
    p = tf.add_paragraph()
    p.text = "Spiritual truth: The goal is not completion. The goal is transformation."
    p = tf.add_paragraph()
    p.text = "💡 Final insight: The Church stretches time so that eternity can enter it"
    p = tf.add_paragraph()
    p.text = "“Each hour is not a service… it is a journey”"

    # 14. What the Music Is Doing to Us
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "6) What the Music Is Doing to Us"
    tf = slide.placeholders[1].text_frame
    tf.text = "Hazayni tone expresses holy sorrow"
    p = tf.add_paragraph()
    p.text = "O Monogenis = creed in suffering"
    p = tf.add_paragraph()
    p.text = "Pekethronos = Christ’s glory in suffering"
    p = tf.add_paragraph()
    p.text = "Repetition acts like spiritual breathing"

    # 15. The Turning Point
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "7) The Turning Point: Where Are You?"
    tf = slide.placeholders[1].text_frame
    tf.text = "Distracted, busy, mechanical, or present?"
    p = tf.add_paragraph()
    p.text = "Am I attending or walking with Christ?"
    
    # 16. Final Movement
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "8) Final Movement: From Length to Love"
    tf = slide.placeholders[1].text_frame
    tf.text = "The Church slows us down to open our hearts"
    p = tf.add_paragraph()
    p.text = "Holy Week is about staying with Christ"

    # 17. Conclusion
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Conclusion"
    tf = slide.placeholders[1].text_frame
    tf.text = "Holy Week transforms sorrow to joy, death to life, ritual to encounter"
    p = tf.add_paragraph()
    p.text = "Every long hour is an invitation: Don’t leave Him"
    
    # Save the file
    prs.save("Holy_Week_Journey.pptx")
    print("PowerPoint generated successfully.")

if __name__ == "__main__":
    create_presentation()
