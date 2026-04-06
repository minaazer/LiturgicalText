import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def apply_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(24, 28, 36)  # Dark elegant blue-grey

def add_title_slide(prs, title_text, subtitle_text, image_path=None):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide)
    
    if image_path:
        # Title slide image positioned centered bottom or background
        try:
            slide.shapes.add_picture(image_path, Inches(3.5), Inches(4.8), height=Inches(2.5))
        except: pass

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0) if image_path else Inches(2.5), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(48)
    p.font.name = "Georgia"
    p.font.bold = True
    p.font.color.rgb = RGBColor(230, 180, 80) # Gold
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.7) if image_path else Inches(4.2), Inches(8), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle_text
    p2.font.size = Pt(28)
    p2.font.name = "Calibri Light"
    p2.font.color.rgb = RGBColor(220, 220, 220)
    p2.alignment = PP_ALIGN.CENTER

    # Decorative Line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.5) if image_path else Inches(4), Inches(2), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(230, 180, 80)
    shape.line.fill.background()
    return slide

def add_content_slide(prs, title_text, points, image_path=None):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_background(slide)
    
    # Slide Title
    txBox_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(8.4), Inches(1))
    tf_title = txBox_title.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(36)
    p_title.font.name = "Georgia"
    p_title.font.color.rgb = RGBColor(230, 180, 80)
    
    # Decorative Divider
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(8.4), Inches(0.02))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(100, 110, 130)
    shape.line.fill.background()

    # Image (if available)
    text_width = Inches(8.4)
    if image_path:
        try:
            # Place image on right
            pic = slide.shapes.add_picture(image_path, Inches(6.0), Inches(2.0), height=Inches(4.8))
            text_width = Inches(5.0)
        except Exception as e:
            print("Could not load image:", e)

    # Content Box
    txBox_content = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), text_width, Inches(4.8))
    tf_content = txBox_content.text_frame
    tf_content.word_wrap = True

    for i, point in enumerate(points):
        text = point.get("text", "")
        level = point.get("level", 0)
        is_insight = point.get("is_insight", False)
        
        if i == 0:
            p = tf_content.paragraphs[0]
        else:
            p = tf_content.add_paragraph()
            
        p.text = "•  " + text if level == 0 and not is_insight else text
        p.level = level
        
        # Space before paragraphs
        p.space_before = Pt(14)
        
        p.font.name = "Calibri Light"
        
        if is_insight:
            p.font.size = Pt(22) if image_path else Pt(24)
            p.font.bold = True
            p.font.italic = True
            p.font.color.rgb = RGBColor(255, 215, 0) # Brighter Gold
        else:
            p.font.size = Pt(22) if image_path and level == 0 else Pt(26) if level == 0 else Pt(18) if image_path else Pt(22)
            p.font.color.rgb = RGBColor(240, 240, 240)

        if level > 0 and not is_insight:
            p.text = "    " + text if level==1 else "        " + text
            
    return slide

def create_presentation():
    import os
    img_dir = r"C:\Users\minae\.gemini\antigravity\brain\6efc4988-9e17-4850-bced-98e025403a1f"
    def get_img(name):
        return os.path.join(img_dir, name)

    prs = Presentation()
    
    add_title_slide(prs, "The Spiritual Journey of Holy Week", "Through Hymns & Rites", get_img("coptic_cross_1774637391763.png"))
    
    add_content_slide(prs, "Entering a Different Kind of Time", [
        {"text": "Holy Week is not a normal week"},
        {"text": "The Church invites us to slow down, step out of ordinary time, and walk with Christ"},
        {"text": "We are not remembering events \u2014 we are walking with Christ"},
        {"text": "Key question: Why is everything so long, repetitive, and musical?", "is_insight": True}
    ], get_img("coptic_light_1774637445060.png"))
    
    add_content_slide(prs, "Why This Week Exists: The Meaning of the Cross", [
        {"text": "Christ entered our pain, fear, and death"},
        {"text": "The Passion is divine love, not tragedy"},
        {"text": "The Cross is where justice meets mercy"}
    ], get_img("coptic_passion_1774637429981.png"))

    add_content_slide(prs, "The Structure of the Journey", [
        {"text": "Palm Sunday \u2192 rejection"},
        {"text": "Monday-Wednesday \u2192 confrontation"},
        {"text": "Thursday \u2192 love and betrayal"},
        {"text": "Friday \u2192 sacrifice"},
        {"text": "Saturday \u2192 waiting"},
        {"text": "Resurrection \u2192 victory"},
        {"text": "Framework: Readings, Rites, Hymns", "is_insight": True}
    ], get_img("coptic_palms_1774637841893.png"))

    add_content_slide(prs, "The Key Triad of Holy Week", [
        {"text": "Readings \u2022 Rites \u2022 Hymns"},
        {"text": "Hymns are theology, prayer, and experience"},
        {"text": "They express what words cannot"},
        {"text": "The Church teaches theology through melody", "is_insight": True}
    ])

    add_content_slide(prs, "Entering One Hour of Holy Week", [
        {"text": "The Prophecies (Old Testament Readings)"},
        {"text": "Drawn from: Isaiah, Jeremiah, Torah, Minor prophets", "level": 1},
        {"text": "Theological meaning: The Cross was not accidental. It was foretold, prepared, and awaited"},
        {"text": "Spiritual purpose: The Church wants you to see Christ everywhere in Scripture and realize that salvation is one continuous story"},
        {"text": "Key insight: The Cross is not a moment\u2026 it is the fulfillment of all history", "is_insight": True}
    ])

    add_content_slide(prs, "Homily", [
        {"text": "The precious teachings of the church fathers that reflect on the events and purpose of this hour or day."},
        {"text": "Liturgical Wisdom: Before the homily the church chants a calm and beautiful tune as an introduction"},
        {"text": "Why is it placed here? Because the readings are long and dense; the mind becomes saturated, tired, full."},
        {"text": "Purpose of the calm melody: it allows digestion of the Word, emotional settling, and preparation for deeper meaning."},
        {"text": "Key insight: The Church understands the human soul and builds rest into worship", "is_insight": True}
    ], get_img("coptic_church_fathers_1774637867892.png"))

    add_content_slide(prs, "Paschal Praise (Thine is the power)", [
        {"text": "Theological meaning: Declares Christ\u2019s kingship and victory"},
        {"text": "Paradox: We are in the middle of suffering, betrayal, crucifixion"},
        {"text": "Yet we say: \u201cYours is the power and the glory\u201d", "level": 1},
        {"text": "Spiritual meaning: The Church sees the Cross not as defeat but as enthronement. The Cross = love + victory, not loss"},
        {"text": "Key insight: We proclaim victory before we see it", "is_insight": True}
    ], get_img("coptic_deacon_1774637881067.png"))

    add_content_slide(prs, "Repetition (Spiritual Formation)", [
        {"text": "The Paschal Praise is repeated many times"},
        {"text": "Why repetition? To move truth from mind \u2192 heart"},
        {"text": "Deeper spiritual role:"},
        {"text": "Creates rhythm, stillness, and internal participation", "level": 1},
        {"text": "Key insight: Repetition turns knowledge into experience. It engraves truth in the heart.", "is_insight": True}
    ])

    add_content_slide(prs, "Psalm in the Long Tune", [
        {"text": "Spiritual meaning: The Psalm is deeply emotional and personal"},
        {"text": "So the Church stretches it, gives space to feel it"},
        {"text": "Effect: Moves from communal reading \u2192 personal prayer"},
        {"text": "Key insight: The Psalm is where you enter the story", "is_insight": True}
    ], get_img("coptic_david_1774637905562.png"))

    add_content_slide(prs, "The Gospel", [
        {"text": "The climax of the hour"},
        {"text": "Meaning: Now Christ speaks directly; the Passion unfolds"},
        {"text": "Spiritual significance: Everything before was preparation. Now comes encounter."},
        {"text": "Connection to structure:"},
        {"text": "Readings \u2192 prepare / Hymns \u2192 soften / Gospel \u2192 reveal", "level": 1},
        {"text": "Key insight: The Gospel is not information. It is presence.", "is_insight": True}
    ], get_img("coptic_hymns_1774637405570.png"))

    add_content_slide(prs, "Exposition (Explanation of the Hour)", [
        {"text": "A poetic, often symbolic explanation"},
        {"text": "Function: Interprets events, symbols and spiritual meaning"},
        {"text": "Style: Not academic; often poetic, symbolic, emotional"},
        {"text": "Spiritual role: Helps the believer locate themselves in the story & understand the deeper meaning"},
        {"text": "Key insight: The exposition teaches us how to see spiritually", "is_insight": True}
    ], get_img("coptic_gospel_1774637917634.png"))

    add_content_slide(prs, "The Hour as a Whole", [
        {"text": "Why is it long? Because each element prepares, deepens, transforms"},
        {"text": "Holy Week is meant for full immersion and devotion"},
        {"text": "Spiritual truth: The goal is not completion. The goal is transformation."},
        {"text": "Final insight: The Church stretches time so that eternity can enter it", "is_insight": True},
        {"text": "\u201cEach hour is not a service\u2026 it is a journey\u201d", "level": 1}
    ])

    add_content_slide(prs, "What the Music Is Doing to Us", [
        {"text": "Hazayni tone expresses holy sorrow"},
        {"text": "O Monogenis = creed in suffering"},
        {"text": "Pekethronos = Christ\u2019s glory in suffering"},
        {"text": "Repetition acts like spiritual breathing"}
    ])

    add_content_slide(prs, "The Turning Point: Where Are You?", [
        {"text": "Distracted, busy, mechanical, or present?"},
        {"text": "Am I attending or walking with Christ?"}
    ])

    add_content_slide(prs, "Final Movement: From Length to Love", [
        {"text": "The Church slows us down to open our hearts"},
        {"text": "Holy Week is about staying with Christ"}
    ])

    add_content_slide(prs, "Conclusion", [
        {"text": "Holy Week transforms sorrow to joy, death to life, ritual to encounter"},
        {"text": "Every long hour is an invitation:"},
        {"text": "\u201cDon\u2019t leave Him\u201d", "is_insight": True}
    ], get_img("coptic_tomb_1774637933293.png"))
    
    # Save the file
    prs.save("Holy_Week_Journey_Beautiful.pptx")
    print("PowerPoint generated successfully.")

if __name__ == "__main__":
    create_presentation()
