import collections
import collections.abc
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = 13.333
SLIDE_H = 7.5

PINK_BG = RGBColor(228, 206, 213)
MAUVE = RGBColor(140, 87, 102)
ROSE = RGBColor(194, 73, 71)
BLUSH = RGBColor(234, 169, 185)
IVORY = RGBColor(250, 245, 226)
CREAM = RGBColor(252, 249, 244)
DARK = RGBColor(79, 64, 66)
BORDER = RGBColor(192, 151, 164)
GOLD = RGBColor(227, 190, 110)

ASSET_DIR = Path("ppt_assets")


def rgb(color):
    return (color[0], color[1], color[2], 255)


def rgba(r, g, b, a=255):
    return (r, g, b, a)


def ensure_assets():
    ASSET_DIR.mkdir(exist_ok=True)
    stack_path = ASSET_DIR / "religious_books_stack.png"
    bible_path = ASSET_DIR / "open_bible_cross.png"
    cross_path = ASSET_DIR / "small_cross.png"
    closed_book_path = ASSET_DIR / "closed_wisdom_book.png"
    collage_path = ASSET_DIR / "wisdom_collage.png"
    if not stack_path.exists():
        create_books_stack_art(stack_path)
    if not bible_path.exists():
        create_open_bible_art(bible_path)
    if not cross_path.exists():
        create_cross_art(cross_path)
    if not closed_book_path.exists():
        create_closed_book_art(closed_book_path)
    if not collage_path.exists():
        create_wisdom_collage_art(collage_path)
    return {
        "stack": str(stack_path.resolve()),
        "bible": str(bible_path.resolve()),
        "cross": str(cross_path.resolve()),
        "closed": str(closed_book_path.resolve()),
        "collage": str(collage_path.resolve()),
    }


def create_books_stack_art(path):
    img = Image.new("RGBA", (1200, 900), (0, 0, 0, 0))
    d = ImageDraw.Draw(img)

    d.ellipse((120, 720, 980, 840), fill=rgba(124, 83, 94, 120))

    def book(x, y, w, h, cover, page):
        d.rounded_rectangle((x, y, x + w, y + h), 45, fill=cover)
        d.rounded_rectangle((x + int(w * 0.55), y + 15, x + w + 30, y + h - 15), 35, fill=page)
        d.rectangle((x + int(w * 0.55) - 12, y + 10, x + int(w * 0.55) + 16, y + h - 10), fill=rgb(MAUVE))
        d.rounded_rectangle((x + 10, y + h - 18, x + w + 20, y + h + 10), 12, fill=rgb(MAUVE))

    book(170, 540, 520, 120, rgba(190, 154, 167), rgba(248, 246, 241))
    book(110, 395, 560, 125, rgba(194, 73, 71), rgba(248, 241, 221))
    book(200, 250, 510, 118, rgba(187, 148, 161), rgba(251, 249, 245))

    d.rounded_rectangle((525, 315, 565, 385), 10, fill=rgba(233, 160, 178))

    d.line((360, 665, 360, 745), fill=rgb(DARK), width=18)
    d.line((320, 705, 400, 705), fill=rgb(DARK), width=18)
    d.line((338, 705, 338, 755), fill=rgb(DARK), width=12)
    d.line((382, 705, 382, 755), fill=rgb(DARK), width=12)

    d.rounded_rectangle((250, 300, 290, 355), 15, fill=rgb(BLUSH))
    d.rounded_rectangle((298, 300, 338, 355), 15, fill=rgb(BLUSH))

    img.save(path)


def create_open_bible_art(path):
    img = Image.new("RGBA", (1200, 900), (0, 0, 0, 0))
    d = ImageDraw.Draw(img)

    d.ellipse((270, 690, 1040, 820), fill=rgba(124, 83, 94, 110))
    d.rounded_rectangle((220, 250, 520, 640), 40, fill=rgba(250, 249, 244))
    d.rounded_rectangle((500, 270, 860, 630), 40, fill=rgba(248, 241, 221))
    d.polygon([(520, 255), (555, 280), (535, 628), (500, 610)], fill=rgba(226, 194, 139))
    d.rectangle((200, 240, 235, 645), fill=rgb(MAUVE))

    for i in range(10):
        y = 320 + i * 27
        d.line((270, y, 465, y), fill=rgba(216, 186, 132), width=3)
        d.line((580, y, 810, y), fill=rgba(216, 186, 132), width=3)

    d.rectangle((810, 180, 1005, 560), fill=rgb(MAUVE))
    d.rectangle((830, 180, 845, 560), fill=rgb(GOLD))
    d.rectangle((860, 250, 960, 285), fill=rgb(ROSE))
    d.rectangle((860, 335, 965, 370), fill=rgb(ROSE))
    d.rectangle((860, 420, 985, 500), fill=rgb(ROSE))

    d.line((898, 610, 898, 740), fill=rgb(DARK), width=20)
    d.line((845, 670, 952, 670), fill=rgb(DARK), width=20)
    d.line((867, 670, 867, 760), fill=rgb(DARK), width=12)
    d.line((929, 670, 929, 760), fill=rgb(DARK), width=12)

    img.save(path)


def create_cross_art(path):
    img = Image.new("RGBA", (180, 220), (0, 0, 0, 0))
    d = ImageDraw.Draw(img)
    d.rounded_rectangle((72, 20, 108, 180), 12, fill=rgb(IVORY))
    d.rounded_rectangle((28, 72, 152, 108), 12, fill=rgb(IVORY))
    img.save(path)


def create_closed_book_art(path):
    img = Image.new("RGBA", (700, 700), (0, 0, 0, 0))
    d = ImageDraw.Draw(img)

    d.rounded_rectangle((140, 120, 540, 560), 28, fill=rgba(191, 167, 178))
    d.rounded_rectangle((135, 135, 185, 548), 22, fill=rgb(MAUVE))
    d.rounded_rectangle((182, 160, 230, 530), 18, fill=rgba(235, 176, 191))
    d.rounded_rectangle((305, 205, 445, 395), 8, fill=rgba(246, 241, 223))
    d.rounded_rectangle((332, 155, 470, 182), 8, fill=rgb(IVORY))
    d.rounded_rectangle((135, 535, 545, 585), 20, fill=rgb(ROSE))
    d.rounded_rectangle((150, 545, 525, 575), 16, fill=rgba(250, 248, 243))
    d.polygon([(205, 552), (255, 552), (248, 585), (195, 585)], fill=rgb(ROSE))
    img.save(path)


def create_wisdom_collage_art(path):
    img = Image.new("RGBA", (1200, 1100), (0, 0, 0, 0))
    d = ImageDraw.Draw(img)

    d.rounded_rectangle((170, 140, 470, 620), 40, fill=rgba(251, 250, 246))
    d.rounded_rectangle((450, 175, 780, 610), 40, fill=rgba(248, 241, 224))
    d.polygon([(160, 165), (195, 140), (215, 185), (176, 204)], fill=rgb(ROSE))
    d.polygon([(166, 580), (200, 553), (222, 603), (182, 620)], fill=rgb(ROSE))
    d.rectangle((145, 155, 170, 615), fill=rgb(MAUVE))
    for i in range(9):
        y = 250 + i * 33
        d.line((245, y, 420, y), fill=rgba(223, 192, 140), width=3)
    d.rectangle((235, 190, 315, 285), outline=rgba(223, 192, 140), width=4)

    d.polygon([(790, 235), (1060, 315), (940, 810), (660, 725)], fill=rgb(MAUVE))
    d.rectangle((815, 335, 995, 670), fill=rgb(ROSE))
    d.rectangle((847, 300, 860, 680), fill=rgb(GOLD))
    d.rectangle((900, 410, 1020, 445), fill=rgba(170, 82, 75))
    d.rectangle((890, 520, 1030, 555), fill=rgba(170, 82, 75))
    d.rectangle((868, 625, 1055, 740), fill=rgba(170, 82, 75))
    d.rectangle((835, 360, 1020, 378), fill=rgba(243, 219, 227))
    d.rectangle((835, 748, 980, 766), fill=rgba(243, 219, 227))

    for cx in (510, 840):
        d.ellipse((cx, 740, cx + 210, 950), fill=rgba(248, 190, 198))
        d.ellipse((cx - 10, 730, cx + 220, 960), outline=rgb(DARK), width=10)
    d.arc((690, 810, 790, 880), 180, 360, fill=rgb(DARK), width=10)
    d.line((500, 875, 445, 915), fill=rgb(DARK), width=18)
    d.line((1055, 875, 1110, 925), fill=rgb(DARK), width=18)
    d.line((495, 875, 470, 835), fill=rgb(DARK), width=12)
    d.line((1060, 875, 1085, 835), fill=rgb(DARK), width=12)
    d.line((562, 775, 650, 855), fill=rgba(255, 255, 255, 210), width=18)
    d.line((625, 755, 708, 835), fill=rgba(255, 255, 255, 210), width=18)
    d.line((892, 775, 980, 855), fill=rgba(255, 255, 255, 210), width=18)
    d.line((955, 755, 1038, 835), fill=rgba(255, 255, 255, 210), width=18)
    img.save(path)


def add_shape(slide, shape_type, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    font_name="Aptos",
    font_size=24,
    color=DARK,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_VERTICAL_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    return tf


def add_border(slide):
    border = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.32),
        Inches(0.28),
        Inches(12.7),
        Inches(6.92),
    )
    border.fill.background()
    border.line.color.rgb = BORDER
    border.line.width = Pt(2.25)
    border.adjustments[0] = 0.08


def apply_background(slide, assets):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PINK_BG
    add_border(slide)
    slide.shapes.add_picture(assets["cross"], Inches(1.1), Inches(1.35), height=Inches(0.72))
    slide.shapes.add_picture(assets["cross"], Inches(11.45), Inches(0.82), height=Inches(0.5))
    slide.shapes.add_picture(assets["cross"], Inches(6.55), Inches(5.55), height=Inches(0.48))


def add_quote_box(slide, left, top, width, height, text):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, RGBColor(234, 217, 223), BORDER)
    tf = add_text(slide, left + 0.22, top + 0.18, width - 0.44, height - 0.3, text, font_size=16, color=MAUVE, bold=True)
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.LEFT


def add_title_slide(prs, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, assets)
    slide.shapes.add_picture(assets["collage"], Inches(0.72), Inches(2.75), width=Inches(5.4))

    add_text(slide, 7.05, 2.55, 4.8, 0.9, "WISDOM", font_name="Georgia", font_size=50, color=MAUVE, bold=True)
    add_text(
        slide, 6.2, 4.05, 6.2, 1.1, "Lesson 5.3: Solomon's Wisdom",
        font_name="Georgia", font_size=28, color=DARK, bold=True, align=PP_ALIGN.CENTER
    )
    add_text(
        slide, 6.3, 5.0, 6.0, 0.9, "God gives wisdom to those who ask Him.",
        font_size=22, color=MAUVE, bold=True, align=PP_ALIGN.CENTER
    )


def add_section_slide(prs, assets, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, assets)
    add_text(slide, 0.92, 0.78, 3.0, 0.65, "WISDOM", font_name="Georgia", font_size=34, color=MAUVE, bold=True)
    add_text(slide, 1.55, 2.0, 4.9, 2.2, title, font_size=30, color=DARK, bold=True)
    add_text(slide, 1.6, 4.95, 4.8, 0.9, subtitle, font_size=22, color=MAUVE)
    slide.shapes.add_picture(assets["bible"], Inches(7.1), Inches(1.45), width=Inches(5.2))


def add_content_slide(prs, assets, title, bullets, art="stack", quote=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, assets)

    add_text(slide, 0.72, 0.62, 3.6, 0.55, "Wisdom", font_name="Georgia", font_size=28, color=MAUVE, bold=True)
    add_text(slide, 0.72, 1.22, 5.9, 0.78, title, font_name="Georgia", font_size=25, color=DARK, bold=True)

    y = 2.05
    for idx, bullet in enumerate(bullets):
        tf = add_text(slide, 1.0, y, 5.7, 0.72, f"• {bullet}", font_size=20, color=DARK if idx < 3 else MAUVE, bold=idx == len(bullets) - 1)
        for p in tf.paragraphs:
            p.space_after = Pt(3)
        y += 0.84

    if quote:
        add_quote_box(slide, 0.86, 5.92, 5.95, 0.72, quote)

    if art == "stack":
        slide.shapes.add_picture(assets["stack"], Inches(7.1), Inches(2.2), width=Inches(5.0))
    elif art == "closed":
        slide.shapes.add_picture(assets["closed"], Inches(8.15), Inches(2.0), width=Inches(3.35))
        slide.shapes.add_picture(assets["cross"], Inches(11.55), Inches(2.45), height=Inches(0.65))
    elif art == "collage":
        slide.shapes.add_picture(assets["collage"], Inches(7.0), Inches(1.95), width=Inches(5.15))
    else:
        slide.shapes.add_picture(assets["bible"], Inches(7.0), Inches(1.85), width=Inches(5.2))


def build_presentation():
    assets = ensure_assets()

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    add_title_slide(prs, assets)

    add_content_slide(
        prs, assets, "Lesson Aim",
        [
            "To understand that God's wisdom is available to everyone who asks Him sincerely.",
            "To see that living with godly wisdom brings blessing, peace, and right judgment.",
            "To learn from Solomon's strengths and from the mistakes that led to his decline.",
            "Wisdom is a gift from God, but it must be guarded with humility.",
        ],
        art="collage",
        quote="Big idea: wisdom is not only knowledge. It is living in a way that pleases God.",
    )

    add_content_slide(
        prs, assets, "Who Was King Solomon?",
        [
            "Solomon was the son of King David and Bathsheba and became the third king of united Israel.",
            "He inherited a strong kingdom and ruled during a time of peace, wealth, and honor.",
            "Under Solomon, Israel reached its greatest point of power, beauty, and influence.",
            "He is remembered as the king most closely connected with wisdom.",
        ],
        art="closed",
        quote="Solomon began with great gifts, a great kingdom, and a great calling.",
    )

    add_content_slide(
        prs, assets, "Why Was Solomon So Famous?",
        [
            "He built the first Temple in Jerusalem, one of the great wonders of the ancient world.",
            "He was known for leadership, prosperity, trade, and peace with surrounding nations.",
            "He wrote proverbs, songs, and books of wisdom that still guide believers today.",
            "His greatness was not only political. It was also spiritual and intellectual.",
        ],
        art="collage",
        quote="God gave Solomon gifts that made his kingdom shine before the nations.",
    )

    add_section_slide(prs, assets, "Solomon Asked\nfor Wisdom", "He asked for the right thing when God tested his heart.")

    add_content_slide(
        prs, assets, "A Wise Request",
        [
            "God appeared to Solomon and told him to ask for whatever he wanted.",
            "Instead of riches or victory over enemies, Solomon asked for understanding.",
            "He wanted wisdom so he could govern God's people with justice and discernment.",
            "Because he asked well, God also gave him honor, riches, and blessing.",
        ],
        art="closed",
        quote="1 Kings 3 shows that God delights in hearts that desire wisdom over selfish gain.",
    )

    add_content_slide(
        prs, assets, "Wisdom on Display",
        [
            "Solomon's wisdom became clear in the famous judgment between the two women and the infant child.",
            "He understood that the true mother would rather give up the child than see him harmed.",
            "His wisdom spread far beyond Israel and reached the ears of the Queen of Sheba.",
            "Real wisdom sees beneath words and reaches the truth of the heart.",
        ],
        art="bible",
        quote="Wisdom is more than information. It is the ability to judge rightly.",
    )

    add_section_slide(prs, assets, "God's Wisdom and\nMan's Wisdom", "True wisdom begins with knowing God and staying humble.")

    add_content_slide(
        prs, assets, "What Is True Wisdom?",
        [
            "Man's wisdom often focuses on intelligence, success, skill, and reputation.",
            "God's wisdom teaches us how to live in a way that pleases Him.",
            "Godly wisdom brings peace, humility, self-control, and right choices.",
            "A proud heart can be smart, but a humble heart can become truly wise.",
        ],
        art="collage",
        quote="Proverbs teaches that wisdom and humility belong together.",
    )

    add_content_slide(
        prs, assets, "What the Fathers Teach Us",
        [
            "Abba Anthony taught that we should always keep God before our eyes.",
            "Abba Daniel warned that too much comfort can weaken the soul.",
            "St. Jerome reminds us that earthly things seem empty when compared with Christ.",
            "The saints show us that wisdom grows best in humility, prayer, and simplicity.",
        ],
        art="closed",
        quote="The Fathers connect wisdom with holiness, not just with cleverness.",
    )

    add_section_slide(prs, assets, "Solomon's Failures", "A gifted life can still fall when the heart drifts from God.")

    add_content_slide(
        prs, assets, "How Solomon Began to Fall",
        [
            "Solomon slowly turned from the pursuit of God toward luxury, pleasure, and comfort.",
            "He allowed wrong influences into his life through disobedience and compromise.",
            "He permitted idolatry and let his heart be divided instead of fully faithful to God.",
            "His story warns us that blessing without watchfulness can become dangerous.",
        ],
        art="collage",
        quote="The greater the gift, the more carefully we must guard the heart.",
    )

    add_content_slide(
        prs, assets, "A Warning for Us",
        [
            "Authority, talent, status, and success do not automatically produce holiness.",
            "When people prosper, they can begin to forget how much they still need God.",
            "We must guard our hearts, our life, and our doctrine with honesty and prayer.",
            "Do not only pray for success. Pray for the character to handle success well.",
        ],
        art="closed",
        quote="Solomon's failure teaches us to stay dependent on God even in blessing.",
    )

    add_content_slide(
        prs, assets, "What Solomon Learned Later",
        [
            "In Ecclesiastes, Solomon admits that a life centered on earthly gain is vanity.",
            "He discovered that pleasure, wealth, and achievement cannot satisfy the heart without God.",
            "His final lesson is simple: fear God and keep His commandments.",
            "The end of wisdom is not pride. It is repentance and obedience.",
        ],
        art="bible",
        quote="Ecclesiastes shows that life without God feels empty no matter how successful it looks.",
    )

    add_section_slide(prs, assets, "How Do We Gain\nWisdom?", "Godly wisdom can grow in us through daily habits of grace.")

    add_content_slide(
        prs, assets, "Ways to Grow in Wisdom",
        [
            "Ask God for wisdom in prayer and speak to Him honestly about your choices.",
            "Read and study the Word of God so your mind learns His ways.",
            "Listen to wise and godly people who can correct and guide you.",
            "Use your gifts faithfully and practice obedience in daily life.",
        ],
        art="closed",
        quote="Prayer, Scripture, humility, and obedience form a wise heart.",
    )

    add_content_slide(
        prs, assets, "Challenge for This Week",
        [
            "Read one chapter of Proverbs each day as a way of training your heart in wisdom.",
            "Pray each day: Lord, give me wisdom and a humble heart.",
            "Choose one area where you will obey God more carefully this week.",
            "Wisdom grows when we not only hear God's word, but also live it.",
        ],
        art="collage",
        quote="Proverbs is one of the clearest daily paths toward wisdom.",
    )

    add_content_slide(
        prs, assets, "Final Thought",
        [
            "Solomon had wisdom, power, honor, and success, yet he still needed faithfulness.",
            "His life teaches us that the greatest treasure is not talent, but closeness to God.",
            "God is still ready to give wisdom generously to those who ask Him.",
            "Wisdom is not only what you know. It is how you walk with God.",
        ],
        art="bible",
        quote='Proverbs 4:11: "I have taught you in the way of wisdom; I have led you in right paths."',
    )

    output_path = "Lesson_5_3_Solomons_Wisdom_Religious_Wisdom_Style_v2.pptx"
    prs.save(output_path)
    print(output_path)


if __name__ == "__main__":
    build_presentation()
